#!/usr/bin/env python3
"""深浅代码风主题 PPT 生成器 — 适配你的答辩内容"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ===== 主题色：深浅代码风 theme03 =====
BG_DARK  = RGBColor(0x12, 0x16, 0x1c)   # 终端深色背景
FG_TEXT  = RGBColor(0xc9, 0xd1, 0xd9)   # 浅灰正文
FG_TITLE = RGBColor(0x58, 0xa6, 0xff)   # 荧光蓝标题
FG_ACCENT= RGBColor(0x23, 0x86, 0xfa)   # 亮蓝强调
FG_GREEN = RGBColor(0x23, 0x86, 0x61)   # 代码绿
FG_WHITE = RGBColor(0xff, 0xff, 0xff)
BORDER   = RGBColor(0x1f, 0x29, 0x37)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_DARK
    return s

def rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = color
        shp.line.width = Pt(0.5)
    return shp

def text(slide, x, y, w, h, content, size=14, bold=False, color=FG_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Microsoft YaHei"
    return tb

def title_bar(slide, eyebrow, main_title):
    # 顶部代码终端蓝条
    rect(slide, 0, 0, SW, Inches(1.0), FG_ACCENT)
    text(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.64), eyebrow, size=16, color=FG_WHITE)
    text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.45), main_title, size=28, bold=True, color=FG_WHITE)
    # 终端装饰小圆点
    for dx, dot_color in [(Inches(0.35), RGBColor(0xff,0x5f,0x57)),
                          (Inches(0.52), RGBColor(0xff,0xbd,0x2d)),
                          (Inches(0.69), RGBColor(0x23,0xc8,0x3b))]:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, dx, Inches(0.3), Inches(0.15), Inches(0.15))
        d.fill.solid()
        d.fill.fore_color.rgb = dot_color
        d.line.fill.background()

# ========== 第1页 封面 ==========
s = new_slide()
title_bar(s, "~/diagnosis_agent", "车辆故障诊断智能助手")
text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.7), "LLM + 双路 RAG + 知识图谱闭环", size=22, color=FG_TITLE)
text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4), "项目名称: Diagnosis Agent v0.5.0", size=14, color=FG_TEXT)
text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4), "技术路线: LangChain ReAct + Neo4j + ChromaDB", size=13, color=RGBColor(0x7d,0x90,0xa3))

# ========== 第2页 目录 ==========
s = new_slide()
title_bar(s, "agenda", "目 录")
chapters = ["01 解决什么问题", "02 技术方案: LLM + 双路 RAG", "03 LangChain 架构", "04 端到端工作流", "05 六大创新点", "06 效果展示数据", "07 技术栈与演示命令"]
for i, c in enumerate(chapters):
    y = Inches(1.6 + i*0.68)
    rect(s, Inches(0.7), y, Inches(0.55), Inches(0.45), FG_GREEN)
    text(s, Inches(0.7), y, Inches(0.55), Inches(0.45), str(i+1), size=20, bold=True, color=FG_WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.5), y+Inches(0.05), Inches(11), Inches(0.4), c, size=17, bold=True, color=FG_TEXT)

# ========== 第3页 三大痛点 ==========
s = new_slide()
title_bar(s, "problem", "经验难沉淀 回答靠幻觉")
pains = [
    ("历史工单散落 Excel", "新人翻几十页找相似案例，效率低"),
    ("诊断报告非结构化", "自然语言写出来很难自动入库统计"),
    ("LLM 凭空回答容易错", "根因对策没有真实证据支撑"),
    ("目标: 基于真实案例推理", "不让大模型凭记忆瞎编")
]
for i, (t, d) in enumerate(pains):
    y = Inches(1.6 + i*1.25)
    rect(s, Inches(0.7), y, Inches(0.6), Inches(0.9), RGBColor(0x1a,0x2f,0x44))
    text(s, Inches(0.7), y, Inches(0.6), Inches(0.9), str(i+1), size=32, bold=True, color=FG_TITLE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.7), y+Inches(0.1), Inches(11), Inches(0.4), t, size=18, bold=True, color=FG_WHITE)
    text(s, Inches(1.7), y+Inches(0.55), Inches(11), Inches(0.35), d, size=13, color=FG_TEXT)

# ========== 第4页 双路 RAG ==========
s = new_slide()
title_bar(s, "tech", "双路 RAG 互补检索")
# 左右卡片
for sx, name, desc_list in [
    (Inches(0.5), "路径1 ChromaDB 向量检索", ["• 语义模糊匹配现象描述", "• 找「行驶中顿挫」相似文本", "• 628 条历史工单", "• 适配: 现象模糊不知道DTC"]),
    (Inches(7.0), "路径2 Neo4j 知识图谱", ["• 按字段精确匹配DTC/车型/场景", "• 2036 节点 / 4356 关系", "• 字段分通道独立匹配零拼接", "• 适配: 有确定字段时零误召"])
]:
    rect(s, sx, Inches(1.6), Inches(5.9), Inches(5.2), RGBColor(0x14,0x1d,0x28))
    rect(s, sx, Inches(1.6), Inches(0.07), Inches(5.2), FG_ACCENT)
    text(s, sx+Inches(0.3), Inches(1.8), Inches(5.5), Inches(0.5), name, size=17, bold=True, color=FG_TITLE)
    for j, line in enumerate(desc_list):
        text(s, sx+Inches(0.3), Inches(2.5 + j*0.65), Inches(5.3), Inches(0.5), line, size=14, color=FG_TEXT)
text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45), "合并去重 → Embedding 精排 → 传给 LLM", size=17, bold=True, color=FG_GREEN, align=PP_ALIGN.CENTER)

# ========== 第5页 六层架构 ==========
s = new_slide()
title_bar(s, "arch", "LangChain 六层分层架构")
layers = [
    ("adapter/", "FastAPI 对接上游平台 异步提交回调"),
    ("agent/", "ReAct 推理循环 编排工具调用"),
    ("retrieval/", "HybridRetriever 双路召回精排"),
    ("knowledge/", "抽取实体关系 审核写入图谱"),
    ("context/", "三层记忆 话题检测 异步摘要"),
    ("storage/", "Neo4j + ChromaDB + Redis + 磁盘")
]
for i, (name, desc) in enumerate(layers):
    y = Inches(1.5 + i*0.86)
    rect(s, Inches(0.7), y, Inches(3.3), Inches(0.72), RGBColor(0x1a,0x2f,0x44))
    text(s, Inches(0.75), y+Inches(0.1), Inches(3.1), Inches(0.52), name, size=15, bold=True, color=FG_TITLE, align=PP_ALIGN.CENTER)
    text(s, Inches(4.2), y+Inches(0.18), Inches(8.5), Inches(0.5), desc, size=14, color=FG_TEXT)

# ========== 第6页 Agent 工具集 ==========
s = new_slide()
title_bar(s, "tools", "动态注册 3 个工具")
table_rows = [
    ["工具", "触发场景", "数据源"],
    ["search_similar_incidents", "模糊匹配现象", "ChromaDB 向量检索"],
    ["query_fault_graph", "结构化字段精确查", "Neo4j Cypher 查询"],
    ["get_incident_detail", "查看工单详情", "历史记录"]
]
for ri, row in enumerate(table_rows):
    y = Inches(2.0 + ri*0.9)
    fillc = RGBColor(0x22,0x30,0x42) if ri == 0 else RGBColor(0x16,0x1f,0x2a)
    xi = [Inches(0.6), Inches(3.8), Inches(9.2)]
    wi = [Inches(3.0), Inches(5.2), Inches(3.5)]
    for ci, cell in enumerate(row):
        rect(s, xi[ci], y, wi[ci], Inches(0.8), fillc)
        tc = FG_WHITE if ri == 0 else FG_TEXT
        b = True if ri == 0 else False
        text(s, xi[ci]+Inches(0.25), y+Inches(0.12), wi[ci]-Inches(0.5), Inches(0.56), cell, size=14, bold=b, color=tc)

# ========== 第7页 端到端流程 ==========
s = new_slide()
title_bar(s, "flow", "完整诊断端到端链路")
steps = [
    ("1", "用户输入", "MCU报P1A3E98爬坡IGBT过温车辆抖动"),
    ("2", "话题判别", "关联上下文排除非诊断话题"),
    ("3", "双路预检索", "Chroma 5条 + Neo4j 3条 加权精排"),
    ("4", "ReAct 推理", "Thought→Action→Obs→Final"),
    ("5", "双层输出", "Markdown人读 + JSON机读入库"),
    ("6", "知识沉淀", "抽取实体关系 写入待审核队列")
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(1.3 + i*0.9)
    rect(s, Inches(0.6), y, Inches(0.6), Inches(0.7), FG_GREEN)
    text(s, Inches(0.6), y, Inches(0.6), Inches(0.7), n, size=22, bold=True, color=FG_WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.6), y+Inches(0.03), Inches(3.5), Inches(0.4), t, size=16, bold=True, color=FG_TITLE)
    text(s, Inches(5.4), y+Inches(0.06), Inches(7.4), Inches(0.55), d, size=13, color=FG_TEXT)

# ========== 第8页 CAN兜底 + 三层记忆 ==========
s = new_slide()
title_bar(s, "mem", "CAN 报文兜底 + 三层记忆")
# 左 CAN
rect(s, Inches(0.5), Inches(1.6), Inches(5.7), Inches(5.2), RGBColor(0x14,0x1d,0x28))
text(s, Inches(0.7), Inches(1.85), Inches(5.3), Inches(0.5), "CAN 报文自动兜底", size=18, bold=True, color=FG_TITLE)
lines_can = [
    "预检索相似度 < 0.6",
    "↓ 解码 CAN 工况文件",
    "支持 ASC/BLF/MF4/CSV",
    "↓ DBC 解码为物理值",
    "↓ 信号摘要注入上下文"
]
for j, l in enumerate(lines_can):
    text(s, Inches(0.9), Inches(2.5 + j*0.55), Inches(5.0), Inches(0.5), l, size=15, color=FG_TEXT)
# 右 三层记忆
rect(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.2), RGBColor(0x14,0x1d,0x28))
text(s, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.5), "三层记忆系统", size=18, bold=True, color=FG_TITLE)
for i, (l, m) in enumerate([("热层 完整消息", "Redis"), ("温层 滚动摘要", "Redis"), ("冷层 归档持久化", "磁盘")]):
    y = Inches(2.5 + i*1.25)
    text(s, Inches(7.4), y, Inches(5), Inches(0.5), f"{l} → {m}", size=16, color=FG_TEXT)

# ========== 第9页 六大创新点 ==========
s = new_slide()
title_bar(s, "feat", "六大核心创新")
innovs = [
    ("双路混合 RAG", "普通 RAG 仅单路向量", "Chroma语义+Neo4j结构化"),
    ("ReAct 推理编排", "普通 LLM 仅生成", "自主调用工具查资料"),
    ("动态知识闭环", "静态知识库", "自动抽取 审核 写入图谱"),
    ("三层分层记忆", "简单多轮对话", "热温冷三级降级"),
    ("CAN 报文自动解码", "无兜底机制", "检索不足注入信号上下文"),
    ("双层输出", "纯文本", "Markdown人读+CSV/JSON机读")
]
for i, (a, b, c) in enumerate(innovs):
    y = Inches(1.5 + i*0.82)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), RGBColor(0x16,0x1f,0x2a))
    text(s, Inches(0.7), y+Inches(0.15), Inches(3.3), Inches(0.5), a, size=15, bold=True, color=FG_TITLE)
    text(s, Inches(4.2), y+Inches(0.18), Inches(3.3), Inches(0.5), b, size=13, color=RGBColor(0x7d,0x90,0xa3))
    text(s, Inches(8.0), y+Inches(0.18), Inches(4.5), Inches(0.5), c, size=14, bold=True, color=FG_TEXT)

# ========== 第10页 效果展示 ==========
s = new_slide()
title_bar(s, "result", "完整诊断输出样例")
text(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.5), "Markdown 人读报告", size=16, bold=True, color=FG_TITLE)
y_off = 2.3
for line in [
    "• 故障分类: 过温故障",
    "• 根因: IGBT 散热基板焊接不良",
    "  动态行驶螺丝松动与PCB短路",
    "• 解决方案: 优化作业流程",
    "• 相似工单 5 条",
    "• 仪表灯: 电机故障红灯"
]:
    text(s, Inches(0.7), y_off, Inches(5.3), Inches(0.4), line, size=14, color=FG_TEXT)
    y_off += 0.55
# JSON 代码块
rect(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0x0f,0x13,0x18))
text(s, Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.4), "JSON 机读结构化输出", size=16, bold=True, color=FG_TITLE)
code_lines = [
    '{',
    '  "classification": "过温故障",',
    '  "root_cause": "IGBT散热基板焊接不良",',
    '  "solution": "优化作业流程",',
    '  "dtc_code": "P1A3E98",',
    '  "confidence": 0.91',
    '}'
]
for j, cl in enumerate(code_lines):
    text(s, Inches(7.25), Inches(2.4 + j*0.52), Inches(5.3), Inches(0.45), cl, size=13, color=RGBColor(0x23,0xe0,0x68))

# ========== 第11页 数据指标 ==========
s = new_slide()
title_bar(s, "metrics", "数据规模与实测")
cards = [
    ("87", "测试用例", "全部通过 41 秒"),
    ("628", "向量库记录", "1024 维 text-embedding-v4"),
    ("2036", "知识图谱节点", "4356 关系 / 8 类实体"),
    ("8000", "Token 预算", "自适应窗口 2~20 轮")
]
for i, (num, label, sub) in enumerate(cards):
    sx = Inches(0.45 + i * 3.25)
    rect(s, sx, Inches(1.7), Inches(3.05), Inches(4.7), RGBColor(0x14,0x1d,0x28))
    rect(s, sx, Inches(1.7), Inches(0.07), Inches(4.7), FG_ACCENT)
    text(s, sx+Inches(0.3), Inches(2.2), Inches(2.5), Inches(1.6), num, size=48, bold=True, color=FG_TITLE, align=PP_ALIGN.CENTER)
    text(s, sx+Inches(0.3), Inches(3.9), Inches(2.5), Inches(0.6), label, size=17, bold=True, color=FG_WHITE, align=PP_ALIGN.CENTER)
    text(s, sx+Inches(0.3), Inches(4.55), Inches(2.5), Inches(0.8), sub, size=12, color=FG_TEXT, align=PP_ALIGN.CENTER)

# ========== 第12页 技术栈 ==========
s = new_slide()
title_bar(s, "stack", "全栈技术选型")
techs = [
    ("LLM 框架", "LangChain ReAct create_agent"),
    ("LLM 后端", "阿里云 DashScope 通义千问"),
    ("知识图谱", "Neo4j 5.22.0 Cypher + APOC"),
    ("向量数据库", "ChromaDB cosine 空间"),
    ("会话存储", "Redis 热温层 + 磁盘冷层归档"),
    ("Web 框架", "FastAPI + Uvicorn"),
    ("CLI", "Typer + Rich 交互式终端"),
    ("语言", "Python 3.10+"),
]
for i, (cat, name) in enumerate(techs):
    y = Inches(1.3 + i*0.72)
    rect(s, Inches(0.6), y, Inches(12.2), Inches(0.6), RGBColor(0x16,0x1f,0x2a))
    text(s, Inches(0.8), y+Inches(0.08), Inches(3.2), Inches(0.45), cat, size=15, bold=True, color=FG_TITLE)
    text(s, Inches(4.3), y+Inches(0.08), Inches(7.8), Inches(0.45), name, size=15, color=FG_WHITE)

# ========== 第13页 演示命令 ==========
s = new_slide()
title_bar(s, "cli", "演示命令")
cmds = [
    ("交互式多轮诊断", "python -m diagnosis_agent.cli chat"),
    ("单次文本诊断", "python -m diagnosis_agent.cli diagnose --text ..."),
    ("加载 CSV 工单", "python -m diagnosis_agent.cli load-data --file ..."),
    ("标准接口对接平台", "python -m diagnosis_agent.cli diagnose --json-input ..."),
    ("查看知识库统计", "python -m diagnosis_agent.cli stats"),
    ("启动 FastAPI 服务", "python -m diagnosis_agent.cli adapter"),
]
for i, (desc, cmd) in enumerate(cmds):
    y = Inches(1.4 + i*0.82)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), RGBColor(0x0f,0x13,0x18))
    text(s, Inches(0.75), y+Inches(0.12), Inches(4.0), Inches(0.46), f"# {desc}", size=13, color=FG_GREEN)
    text(s, Inches(5.0), y+Inches(0.12), Inches(7.3), Inches(0.46), cmd, size=14, color=RGBColor(0x8b,0xc3,0x4a))

# ========== 第14页 结束页 ==========
s = new_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = BG_DARK
text(s, 0, Inches(2.4), SW, Inches(1.4), "谢谢", size=72, bold=True, color=FG_TITLE, align=PP_ALIGN.CENTER)
text(s, 0, Inches(4.1), SW, Inches(0.8), "欢迎各位老师提问", size=26, color=FG_TEXT, align=PP_ALIGN.CENTER)

out_path = "/home/dfmc/diagnose_agent/shared/diagnosis_agent/docs/答辩文档_theme03.pptx"
prs.save(out_path)
print(f"✅ 已生成深浅代码风主题 PPT: {out_path}")
print(f"  总页数: {len(prs.slides)}")
