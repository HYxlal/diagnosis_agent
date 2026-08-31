#!/usr/bin/env python3
"""答辩 PPT 最终版 — 商务浅蓝配色，第一版原始布局 + 轻微美化，不花哨"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== 商务浅蓝配色 — 学术答辩风格 =====
C_BG     = RGBColor(0xf6, 0xf8, 0xfa)   # 浅灰蓝背景
C_BG2    = RGBColor(0xeb, 0xf2, 0xf9)   # 卡片浅蓝背景
C_PRIMARY= RGBColor(0x0f, 0x4c, 0x8c)   # 主色深蓝
C_ACCENT = RGBColor(0x25, 0x63, 0xeb)   # 强调亮蓝
C_GREEN  = RGBColor(0x0e, 0x9f, 0x5f)   # 点缀绿
C_TEXT   = RGBColor(0x1a, 0x20, 0x2c)   # 正文深灰黑
C_SUB    = RGBColor(0x4b, 0x55, 0x65)   # 次要文字灰
C_LINE   = RGBColor(0xd0, 0xd8, 0xe3)   # 分割线浅灰蓝
C_WHITE  = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def fill_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rounded(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = C_LINE
    shp.line.width = Pt(0.5)
    try:
        shp.adjustments[0] = 0.05
    except:
        pass
    return shp

def text(slide, x, y, w, h, content, size=14, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return tb

def title_bar(slide, eyebrow, main_title, page_no, total):
    """第一版原始头部 + 美化"""
    # 顶部主色条
    rect(slide, 0, 0, SW, Inches(0.9), C_PRIMARY)
    text(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.5), eyebrow, size=13, color=RGBColor(0xb8, 0xd4, 0xf1))
    text(slide, Inches(0.5), Inches(0.42), Inches(10), Inches(0.45), main_title, size=24, bold=True, color=C_WHITE)
    # 页码
    text(slide, Inches(12.0), Inches(0.42), Inches(1.0), Inches(0.45), f"{page_no}/{total}", size=12, color=RGBColor(0xb8, 0xd4, 0xf1), align=PP_ALIGN.RIGHT)
    # 底部细蓝条
    rect(slide, 0, Inches(7.3), SW, Inches(0.2), C_BG2)

# ========== 第1页 封面 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
# 主色背景
s.background.fill.solid()
s.background.fill.fore_color.rgb = C_PRIMARY

text(s, Inches(1.0), Inches(1.8), Inches(12), Inches(1.1), "车辆故障诊断智能助手", size=52, bold=True, color=C_WHITE)
text(s, Inches(1.0), Inches(3.0), Inches(12), Inches(0.6), "LLM + 双路 RAG + 知识图谱闭环", size=24, color=RGBColor(0xb0, 0xd3, 0xf5))
text(s, Inches(1.0), Inches(4.2), Inches(12), Inches(0.45), "项目名称：Diagnosis Agent v0.5.0", size=16, color=C_WHITE)
text(s, Inches(1.0), Inches(4.8), Inches(12), Inches(0.45), "技术路线：LangChain ReAct + Neo4j + ChromaDB", size=14, color=RGBColor(0xb8, 0xd4, 0xf1))

# ========== 第2页 目录 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "目录", "目 录", 2, 14)
chapters = [
    ("01", "解决什么问题"), ("02", "技术方案：LLM + 双路 RAG"),
    ("03", "项目搭建思路：LangChain 架构"), ("04", "实际工作流"),
    ("05", "核心创新点"), ("06", "效果展示"),
    ("07", "技术栈"), ("08", "演示命令")
]
for i, (no, name) in enumerate(chapters):
    col = i // 4
    row = i % 4
    x = Inches(0.7 + col * 6.2)
    y = Inches(1.6 + row * 1.3)
    rounded(s, x, y, Inches(0.75), Inches(0.75), C_ACCENT)
    text(s, x, y, Inches(0.75), Inches(0.75), no, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text(s, x + Inches(1.0), y + Inches(0.1), Inches(4.5), Inches(0.6), name, size=17, bold=True, color=C_TEXT)

# ========== 第3页 三大痛点 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "一、解决什么问题？", "传统诊断三大痛点", 3, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "传统车辆故障诊断面临三个痛点，根本上都是「经验难沉淀，回答靠幻觉」", size=14, color=C_SUB)
pains = [
    ("历史工单散落 Excel", "新人翻几十页才能找到相似案例，效率低"),
    ("诊断报告非结构化", "自然语言写出来很难自动入库统计"),
    ("LLM 凭空回答容易错", "不知道真实历史案例，根因对策无证据支撑"),
    ("目标：基于真实案例推理", "不让大模型凭记忆瞎编")
]
for i, (t, d) in enumerate(pains):
    y = Inches(2.0 + i * 1.3)
    rounded(s, Inches(0.7), y, Inches(0.8), Inches(0.8), C_ACCENT)
    text(s, Inches(0.7), y, Inches(0.8), Inches(0.8), str(i+1), size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.8), y+Inches(0.05), Inches(11), Inches(0.5), t, size=19, bold=True, color=C_PRIMARY)
    text(s, Inches(1.8), y+Inches(0.55), Inches(11), Inches(0.4), d, size=14, color=C_TEXT)

# ========== 第4页 双路 RAG ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "二、技术方案", "LLM + 双路 RAG", 4, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "普通 RAG 只有向量检索一条路，我们做了两条互补路径，精度更高", size=14, color=C_SUB)

# 左右两个卡片
for sx, title_text, items_list in [
    (Inches(0.5), "路径1：ChromaDB 向量检索", [
        "语义模糊匹配现象描述",
        "找「行驶中顿挫」等相似文本片段",
        "628 条历史工单，1024 维向量",
        "适配：现象模糊不知道精确 DTC 码时"
    ]),
    (Inches(7.0), "路径2：Neo4j 知识图谱", [
        "按字段精确匹配 DTC/车型/场景",
        "2036 节点 / 4356 关系",
        "字段分通道独立匹配，零拼接",
        "适配：有确定字段时，零误召"
    ])
]:
    rounded(s, sx, Inches(1.85), Inches(5.8), Inches(2.7), C_BG2)
    rect(s, sx, Inches(1.85), Inches(5.8), Inches(0.55), C_PRIMARY)
    text(s, sx+Inches(0.3), Inches(1.92), Inches(5.2), Inches(0.4), title_text, size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items_list):
        text(s, sx+Inches(0.45), Inches(2.6 + j*0.52), Inches(5.0), Inches(0.45), f"• {item}", size=14, color=C_TEXT)

# 底部融合条
rounded(s, Inches(2.2), Inches(4.8), Inches(9.0), Inches(0.75), C_GREEN)
text(s, Inches(2.2), Inches(4.8), Inches(9.0), Inches(0.75),
     "结果合并去重 → Embedding 精排 → 传给 LLM", size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
     "两条路径各自独立匹配，完全不做字段字符串拼接，最后由精排层融合排序。", size=13, color=C_SUB, align=PP_ALIGN.CENTER)

# ========== 第5页 六层分层架构 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "三、项目搭建思路", "LangChain 六层架构", 5, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "LangChain 提供三大抽象，每层职责单一，通过标准接口串联", size=14, color=C_SUB)

layers = [
    ("适配层 adapter/", "FastAPI 对接上游平台（异步提交 + 回调）"),
    ("Agent 层 agent/", "ReAct 推理循环，编排工具调用"),
    ("检索层 retrieval/", "HybridRetriever 双路召回 + 精排"),
    ("知识层 knowledge/", "对话抽取实体关系 → 审核 → 写回图谱"),
    ("上下文层 context/", "三层记忆 + 话题检测 + 异步摘要"),
    ("存储层 storage/", "Neo4j + ChromaDB + Redis + 磁盘")
]
for i, (name, desc) in enumerate(layers):
    y = Inches(1.45 + i*0.83)
    rounded(s, Inches(0.6), y, Inches(12.1), Inches(0.7), C_BG2)
    rect(s, Inches(0.6), y, Inches(0.07), Inches(0.7), C_ACCENT)
    text(s, Inches(0.9), y+Inches(0.12), Inches(3.0), Inches(0.45), name, size=14, bold=True, color=C_PRIMARY)
    text(s, Inches(4.1), y+Inches(0.12), Inches(8.3), Inches(0.45), desc, size=14, color=C_TEXT)

# ========== 第6页 Agent 工具集 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "三、Agent 工具集", "动态注册 3 个工具", 6, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "通过 LangChain StructuredTool 暴露 3 个工具，不可用时自动降级不注册", size=14, color=C_SUB)

rounded(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.1), C_BG2)
# 表头
rect(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6), C_PRIMARY)
text(s, Inches(0.7), Inches(1.98), Inches(3.8), Inches(0.45), "工具", size=14, bold=True, color=C_WHITE)
text(s, Inches(4.7), Inches(1.98), Inches(4.0), Inches(0.45), "触发场景", size=14, bold=True, color=C_WHITE)
text(s, Inches(9.0), Inches(1.98), Inches(3.5), Inches(0.45), "数据源", size=14, bold=True, color=C_WHITE)

tools_data = [
    ("search_similar_incidents", "模糊匹配故障现象", "ChromaDB 向量检索"),
    ("query_fault_graph", "有结构化字段精确查", "Neo4j Cypher 查询"),
    ("get_incident_detail", "查看工单详情", "历史记录")
]
for i, (name, trig, src) in enumerate(tools_data):
    y = Inches(2.55 + i*0.82)
    bg = C_WHITE if i % 2 == 0 else C_BG2
    rounded(s, Inches(0.6), y, Inches(12.1), Inches(0.75), bg)
    text(s, Inches(0.8), y+Inches(0.15), Inches(3.6), Inches(0.45), name, size=13, color=C_PRIMARY, font_name="Consolas")
    text(s, Inches(4.7), y+Inches(0.15), Inches(4.0), Inches(0.45), trig, size=13, color=C_TEXT)
    text(s, Inches(9.0), y+Inches(0.15), Inches(3.3), Inches(0.45), src, size=13, color=C_TEXT)

# ========== 第7页 端到端流程 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "四、实际工作流", "端到端 6 步链路", 7, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "完整诊断从用户提问到输出，经过以下步骤", size=14, color=C_SUB)

steps = [
    ("1", "用户输入", "MCU报P1A3E98爬坡IGBT过温，车辆抖动"),
    ("2", "话题命中判别", "关联上下文，排除非车辆故障诊断话题"),
    ("3", "预检索（双路 RAG 召回）", "Chroma 5条 + Neo4j 3条 → 加权打分精排"),
    ("4", "ReAct Agent 推理循环", "Thought → Action → Observation → Final"),
    ("5", "可选双层输出", "Markdown 人读 + CSV/JSON 机读入库"),
    ("6", "知识沉淀", "LLM 抽取实体/关系 → 写入待审核队列")
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(1.45 + i*0.88)
    rounded(s, Inches(0.6), y, Inches(12.1), Inches(0.72), C_BG2)
    # 序号圆
    rounded(s, Inches(0.75), y+Inches(0.07), Inches(0.58), Inches(0.58), C_ACCENT)
    text(s, Inches(0.75), y+Inches(0.07), Inches(0.58), Inches(0.58), n, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.6), y+Inches(0.08), Inches(3.8), Inches(0.55), t, size=16, bold=True, color=C_PRIMARY)
    text(s, Inches(5.6), y+Inches(0.10), Inches(6.7), Inches(0.5), d, size=13, color=C_TEXT)

# ========== 第8页 CAN 兜底 + 三层记忆 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "四、CAN 报文兜底", "三层记忆系统", 8, 14)

# 左右分栏卡片
rounded(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.3), C_BG2)
text(s, Inches(0.7), Inches(1.65), Inches(5.3), Inches(0.5), "CAN 报文自动兜底", size=17, bold=True, color=C_PRIMARY)
can_lines = [
    "当预检索结果不足（top1 相似度 < 0.6）时：",
    "",
    "预检索相似度低",
    "↓ 解码 CAN 工况文件",
    "支持 ASC/BLF/MF4/CSV 格式",
    "↓ DBC 解码为物理值",
    "↓ 信号摘要注入 Agent 上下文",
    "Agent 基于信号数据推理"
]
for j, l in enumerate(can_lines):
    text(s, Inches(0.9), Inches(2.3 + j*0.48), Inches(5.0), Inches(0.4), l, size=13, color=C_TEXT)

rounded(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.3), C_BG2)
text(s, Inches(7.2), Inches(1.65), Inches(5.3), Inches(0.5), "三层记忆系统", size=17, bold=True, color=C_PRIMARY)
mem_items = [
    ("热层", "完整消息", "Redis"),
    ("温层", "滚动摘要", "Redis"),
    ("冷层", "归档持久化", "磁盘")
]
for i, (n, c, sname) in enumerate(mem_items):
    y = Inches(2.4 + i*1.3)
    rounded(s, Inches(7.4), y, Inches(5.0), Inches(1.05), C_WHITE)
    text(s, Inches(7.6), y+Inches(0.1), Inches(1.5), Inches(0.4), n, size=15, bold=True, color=C_ACCENT)
    text(s, Inches(9.3), y+Inches(0.1), Inches(2.8), Inches(0.45), c, size=13, color=C_TEXT)
    text(s, Inches(11.5), y+Inches(0.25), Inches(1.0), Inches(0.4), sname, size=12, color=C_GREEN, align=PP_ALIGN.CENTER)

# ========== 第9页 六大创新点 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "五、核心创新点", "六大改进远超普通 RAG", 9, 14)
text(s, Inches(0.5), Inches(1.28), Inches(12), Inches(0.4),
     "与普通 RAG 相比的六大创新能力升级", size=14, color=C_SUB)

innovs = [
    ("检索模式", "单路向量", "双路混合 RAG（Chroma 语义 + Neo4j 结构化）"),
    ("LLM 角色", "上下文生成器", "ReAct 推理编排者，主动调用工具查资料"),
    ("知识库", "静态", "动态闭环（自动抽取 → 审核 → 写入图谱）"),
    ("多轮记忆", "无/简单", "三层分层记忆（热/温/冷），支持话题切换"),
    ("兜底机制", "无", "CAN 报文自动解码，检索不足时注入信号上下文"),
    ("输出格式", "纯文本", "双层输出（人读 Markdown + 机读 CSV/JSON）")
]
for i, (a, b, c) in enumerate(innovs):
    bgc = C_WHITE if i % 2 == 0 else C_BG2
    rounded(s, Inches(0.5), Inches(1.7 + i*0.72), Inches(12.3), Inches(0.65), bgc)
    rect(s, Inches(0.5), Inches(1.7 + i*0.72), Inches(0.07), Inches(0.65), C_GREEN)
    text(s, Inches(0.9), Inches(1.82 + i*0.72), Inches(2.8), Inches(0.42), a, size=14, bold=True, color=C_PRIMARY)
    text(s, Inches(4.1), Inches(1.82 + i*0.72), Inches(3.5), Inches(0.42), b, size=13, color=C_SUB)
    text(s, Inches(7.8), Inches(1.82 + i*0.72), Inches(4.7), Inches(0.42), c, size=13, bold=True, color=C_GREEN)

# ========== 第10页 效果展示 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "六、效果展示", "完整诊断输出样例", 10, 14)

# 左右两个卡片
rounded(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4), C_BG2)
text(s, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.45), "人类可读 Markdown 报告", size=16, bold=True, color=C_PRIMARY)
for j, l in enumerate([
    "• 故障分类：过温故障",
    "• 根因分析：IGBT 散热基板焊接不良",
    "  动态行驶螺丝松动与 PCB 电容接触短路",
    "• 解决方案：优化作业流程",
    "  将螺纹孔深度纳入作业流程管理",
    "• 相似工单：5 条（含车型/工况/软件版本）",
    "• 仪表指示灯：电机故障红灯"
]):
    text(s, Inches(0.9), Inches(2.2 + j*0.5), Inches(5.1), Inches(0.42), l, size=14, color=C_TEXT)

rounded(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.4), C_BG2)
text(s, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.45), "机器可读 JSON 条目", size=16, bold=True, color=C_PRIMARY)
for j, l in enumerate([
    '{',
    '  "classification": "过温故障",',
    '  "root_cause": "IGBT散热基板焊接不良",',
    '  "solution": "优化作业流程",',
    '  "dtc_code": "P1A3E98",',
    '  "similar_record_ids": [...]',
    '  "confidence": 0.91',
    '}'
]):
    text(s, Inches(7.4), Inches(2.2 + j*0.42), Inches(5.2), Inches(0.38), l, size=13, color=C_GREEN, font_name="Consolas")
text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
     "故障分类强制从 10 类固定列表选择，保证分类一致性", size=12, color=C_SUB, align=PP_ALIGN.CENTER)

# ========== 第11页 数据规模 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "六、数据规模与测试", "实测指标", 11, 14)

cards_data = [
    ("87", "测试用例", "全部通过，41 秒"),
    ("628", "向量库记录", "1024 维 text-embedding-v4"),
    ("2036", "知识图谱节点", "4356 关系 / 8 类实体"),
    ("8000", "Token 预算", "热层窗口自适应 2~20 轮")
]
for i, (num, label, sub) in enumerate(cards_data):
    x = Inches(0.4 + i * 3.2)
    rounded(s, x, Inches(1.6), Inches(3.05), Inches(3.6), C_BG2)
    text(s, x, Inches(2.2), Inches(3.05), Inches(1.3), num, size=50, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.55), Inches(3.05), Inches(0.55), label, size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.15), Inches(4.2), Inches(2.75), Inches(0.7), sub, size=12, color=C_SUB, align=PP_ALIGN.CENTER)

# ========== 第12页 技术栈 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "七、技术栈", "完整选型清单", 12, 14)

techs = [
    ("LLM 框架", "LangChain create_agent ReAct 模式"),
    ("LLM 后端", "阿里云 DashScope（通义千问，OpenAI 兼容）"),
    ("知识图谱", "Neo4j 5.22.0（Cypher + APOC）"),
    ("向量数据库", "ChromaDB（cosine 空间，PersistentClient）"),
    ("会话存储", "Redis（热层/温层）+ 磁盘（冷层归档）"),
    ("Web 框架", "FastAPI + Uvicorn（平台适配层）"),
    ("CLI", "Typer + Rich（交互式终端）"),
    ("语言", "Python 3.10+")
]
for i, (cat, name) in enumerate(techs):
    y = Inches(1.4 + i*0.72)
    rounded(s, Inches(0.5), y, Inches(12.3), Inches(0.62), C_BG2)
    text(s, Inches(0.8), y+Inches(0.1), Inches(3.0), Inches(0.42), cat, size=14, color=C_SUB)
    text(s, Inches(4.2), y+Inches(0.1), Inches(7.7), Inches(0.42), name, size=15, color=C_TEXT)

# ========== 第13页 演示命令 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
title_bar(s, "八、演示命令", "常用命令列表", 13, 14)

cmds = [
    ("交互式多轮诊断", "python -m diagnosis_agent.cli chat"),
    ("单次文本诊断", "python -m diagnosis_agent.cli diagnose --text ..."),
    ("加载 CSV 工单数据", "python -m diagnosis_agent.cli load-data --file ..."),
    ("标准接口对接平台", "python -m diagnosis_agent.cli diagnose --json-input ..."),
    ("查看知识库统计", "python -m diagnosis_agent.cli stats"),
    ("启动 FastAPI 服务", "python -m diagnosis_agent.cli adapter")
]
for i, (desc, cmd) in enumerate(cmds):
    y = Inches(1.4 + i*0.8)
    rounded(s, Inches(0.5), y, Inches(12.3), Inches(0.68), C_BG2)
    text(s, Inches(0.8), y+Inches(0.1), Inches(4.2), Inches(0.48), f"# {desc}", size=12, color=C_SUB)
    text(s, Inches(5.3), y+Inches(0.1), Inches(7.0), Inches(0.48), cmd, size=13, color=C_GREEN, font_name="Consolas")

# ========== 第14页 结束页 ==========
s = prs.slides.add_slide(BLANK)
fill_bg(s)
s.background.fill.solid()
s.background.fill.fore_color.rgb = C_PRIMARY
text(s, 0, Inches(2.5), SW, Inches(1.1), "谢谢", size=68, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
text(s, 0, Inches(3.8), SW, Inches(0.7), "欢迎提问", size=24, color=RGBColor(0xb8,0xd4,0xf1), align=PP_ALIGN.CENTER)

out_path = "/home/dfmc/diagnose_agent/shared/diagnosis_agent/docs/答辩文档_final.pptx"
prs.save(out_path)
print(f"✅ 已生成最终版: {out_path}")
print(f"  页数: {len(prs.slides)}")
