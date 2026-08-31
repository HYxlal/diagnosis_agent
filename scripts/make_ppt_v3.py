#!/usr/bin/env python3
"""答辩 PPT v3 — 多层次配色 + 图标 + 渐变 + 数据可视化"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ===== 多层次配色（6+ 色系） =====
C_BG1    = RGBColor(0x0d, 0x1b, 0x2a)   # 深海军蓝主背景
C_BG2    = RGBColor(0x1a, 0x2a, 0x3e)   # 次级深蓝
C_BG3    = RGBColor(0x22, 0x33, 0x4a)   # 卡片背景
C_PANEL  = RGBColor(0x2a, 0x3f, 0x5c)   # 面板色
C_CYAN   = RGBColor(0x00, 0xd4, 0xff)   # 青色高亮
C_BLUE   = RGBColor(0x3b, 0x82, 0xf6)   # 亮蓝
C_PURPLE = RGBColor(0x8b, 0x5c, 0xf6)   # 紫色
C_GREEN  = RGBColor(0x10, 0xb9, 0x81)   # 翠绿
C_ORANGE = RGBColor(0xf5, 0x9e, 0x0b)   # 琥珀
C_PINK   = RGBColor(0xec, 0x48, 0x99)   # 粉色
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LIGHT  = RGBColor(0xcb, 0xd5, 0xe1)   # 浅灰白
C_MUTE   = RGBColor(0x64, 0x74, 0x8b)   # 静音灰
C_LINE   = RGBColor(0x33, 0x4d, 0x6e)   # 分割线

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def gradient_bg(slide, color1, color2):
    """渐变背景"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1

def slide_dark(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG1

def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def rounded(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    # 调整圆角半径
    try:
        shp.adjustments[0] = 0.06
    except:
        pass
    return shp

def card(slide, x, y, w, h, bg_color=C_BG3, border_color=None):
    """带阴影的卡片"""
    # 阴影层
    shadow = rect(slide, x + Emu(38100), y + Emu(38100), w, h, C_BG1)
    # 卡片本体
    c = rounded(slide, x, y, w, h, bg_color)
    if border_color:
        c.line.color.rgb = border_color
        c.line.width = Pt(1)
    return c

def text(slide, x, y, w, h, content, size=14, bold=False, color=C_LIGHT, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb

def icon(slide, x, y, size, symbol, bg_color, icon_color, shape=MSO_SHAPE.OVAL):
    """带背景色的图标圆"""
    s = rect(slide, x, y, size, size, bg_color, shape=shape)
    text(slide, x, y, size, size, symbol, size=int(size/9525), bold=True, color=icon_color, align=PP_ALIGN.CENTER)
    return s

def header(slide, section_no, title, accent_color=C_CYAN):
    """统一的页面头部"""
    # 左侧装饰条
    rect(slide, 0, 0, Inches(0.12), SH, accent_color)
    # 顶部小标签
    rounded(slide, Inches(0.4), Inches(0.35), Inches(0.55), Inches(0.4), C_BG2)
    text(slide, Inches(0.4), Inches(0.35), Inches(0.55), Inches(0.4), section_no, size=14, bold=True, color=accent_color, align=PP_ALIGN.CENTER)
    # 主标题
    text(slide, Inches(1.1), Inches(0.3), Inches(11), Inches(0.6), title, size=28, bold=True, color=C_WHITE)
    # 分割线
    rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Pt(1), C_LINE)

def chip(slide, x, y, w, label, color):
    """小标签芯片"""
    s = rounded(slide, x, y, w, Inches(0.32), color)
    s.fill.solid()
    # 半透明效果用深色模拟
    s.fill.fore_color.rgb = RGBColor(color[0]//3, color[1]//3, color[2]//3) if isinstance(color, RGBColor) else color
    text(slide, x, y, w, Inches(0.32), label, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


# ========== 第1页 封面 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
# 背景大色块装饰
rect(s, 0, 0, SW, SH, C_BG1)
# 右侧渐变色块
rect(s, Inches(8.5), 0, Inches(4.833), SH, C_BG2)
# 顶部装饰
rect(s, 0, 0, SW, Inches(0.15), C_CYAN)
rect(s, 0, Inches(0.15), Inches(4), Inches(0.04), C_BLUE)
rect(s, Inches(4.2), Inches(0.15), Inches(3), Inches(0.04), C_PURPLE)

# 图标圆
icon(s, Inches(0.8), Inches(1.2), Inches(1.2), "🚗", C_BG3, C_CYAN)

# 标题
text(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.0), "车辆故障诊断", size=44, bold=True, color=C_WHITE)
text(s, Inches(0.8), Inches(3.4), Inches(7), Inches(1.0), "智能助手", size=44, bold=True, color=C_CYAN)

# 副标题条
rect(s, Inches(0.85), Inches(4.5), Inches(0.04), Inches(0.5), C_GREEN)
text(s, Inches(1.05), Inches(4.5), Inches(7), Inches(0.5), "LLM + 双路 RAG + 知识图谱闭环", size=20, color=C_LIGHT)

# 右侧信息卡
card(s, Inches(8.9), Inches(2.0), Inches(4.0), Inches(3.5), C_BG3, C_CYAN)
text(s, Inches(9.2), Inches(2.3), Inches(3.5), Inches(0.4), "项目信息", size=14, bold=True, color=C_CYAN)
infos = [("项目名称", "Diagnosis Agent v0.5.0"),
         ("技术路线", "LangChain ReAct"),
         ("核心组件", "Neo4j + ChromaDB"),
         ("适用场景", "电驱系统故障诊断")]
for i, (k, v) in enumerate(infos):
    y = Inches(2.85 + i * 0.6)
    text(s, Inches(9.2), y, Inches(1.5), Inches(0.5), k, size=11, color=C_MUTE)
    text(s, Inches(10.7), y, Inches(2.2), Inches(0.5), v, size=12, bold=True, color=C_WHITE)

# ========== 第2页 目录 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "00", "目 录", C_CYAN)

chapters = [
    ("01", "解决什么问题", "三大痛点", C_ORANGE),
    ("02", "技术方案：双路 RAG", "Chroma + Neo4j", C_CYAN),
    ("03", "LangChain 架构", "六层分层", C_BLUE),
    ("04", "端到端工作流", "6 步链路", C_PURPLE),
    ("05", "六大核心创新点", "远超普通 RAG", C_GREEN),
    ("06", "效果展示与数据", "实测结果", C_PINK),
    ("07", "技术栈与演示", "全栈选型", C_ORANGE)
]
for i, (no, title, sub, color) in enumerate(chapters):
    col = i // 4
    row = i % 4
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.5 + row * 1.35)
    # 卡片
    card(s, x, y, Inches(6.0), Inches(1.15), C_BG3)
    # 左侧色条
    rect(s, x, y, Inches(0.08), Inches(1.15), color)
    # 序号圆
    icon(s, x + Inches(0.25), y + Inches(0.2), Inches(0.75), no, color, C_WHITE)
    # 标题
    text(s, x + Inches(1.3), y + Inches(0.15), Inches(4.5), Inches(0.5), title, size=17, bold=True, color=C_WHITE)
    text(s, x + Inches(1.3), y + Inches(0.65), Inches(4.5), Inches(0.35), sub, size=12, color=color)

# ========== 第3页 三大痛点 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "01", "解决什么问题？", C_ORANGE)
text(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
     "传统车辆故障诊断面临三个痛点，根本上都是「经验难沉淀，回答靠幻觉」", size=14, color=C_MUTE)

pains = [
    ("📋", "历史工单散落 Excel", "新人翻几十页才能找到相似案例，效率低下", C_ORANGE),
    ("📝", "诊断报告非结构化", "自然语言写出来很难自动入库统计", C_PINK),
    ("🤖", "LLM 凭空回答容易错", "不知道真实历史案例，根因对策无证据支撑", C_PURPLE)
]
for i, (icon_sym, t, d, color) in enumerate(pains):
    y = Inches(2.0 + i * 1.35)
    card(s, Inches(0.5), y, Inches(8.0), Inches(1.15), C_BG3)
    rect(s, Inches(0.5), y, Inches(0.08), Inches(1.15), color)
    # 图标
    icon(s, Inches(0.8), y + Inches(0.25), Inches(0.65), icon_sym, C_BG2, color)
    text(s, Inches(1.7), y + Inches(0.15), Inches(6.5), Inches(0.5), t, size=17, bold=True, color=C_WHITE)
    text(s, Inches(1.7), y + Inches(0.65), Inches(6.5), Inches(0.4), d, size=13, color=C_LIGHT)

# 右侧目标卡
card(s, Inches(8.8), Inches(2.0), Inches(4.2), Inches(3.85), C_BG2, C_GREEN)
icon(s, Inches(9.1), Inches(2.3), Inches(0.8), "🎯", RGBColor(0x0a, 0x2e, 0x23), C_GREEN)
text(s, Inches(9.1), Inches(3.3), Inches(3.5), Inches(0.5), "我们的目标", size=16, bold=True, color=C_GREEN)
text(s, Inches(9.1), Inches(3.9), Inches(3.5), Inches(1.5),
     "让大模型基于真实历史案例推理而不是凭记忆瞎编", size=15, color=C_WHITE)

# ========== 第4页 双路 RAG ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "02", "技术方案：LLM + 双路 RAG", C_CYAN)
text(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
     "普通 RAG 只有向量检索一条路，我们做了两条互补路径，精度更高", size=14, color=C_MUTE)

# 左卡 ChromaDB
card(s, Inches(0.5), Inches(1.85), Inches(5.8), Inches(4.0), C_BG3, C_CYAN)
rect(s, Inches(0.5), Inches(1.85), Inches(5.8), Inches(0.7), C_BG2)
icon(s, Inches(0.8), Inches(1.95), Inches(0.5), "🔍", C_BG2, C_CYAN)
text(s, Inches(1.5), Inches(1.95), Inches(4.5), Inches(0.5), "路径1: ChromaDB 向量检索", size=16, bold=True, color=C_CYAN)
items_chroma = [
    ("语义模糊匹配现象描述", C_CYAN),
    ("找「行驶中顿挫」相似文本", C_LIGHT),
    ("628 条历史工单 / 1024 维", C_LIGHT),
    ("适配: 现象模糊场景", C_GREEN)
]
for i, (item, c) in enumerate(items_chroma):
    y = Inches(2.75 + i * 0.6)
    rect(s, Inches(0.9), y + Inches(0.08), Inches(0.08), Inches(0.08), c)
    text(s, Inches(1.1), y, Inches(4.8), Inches(0.5), item, size=14, color=c, bold=(c==C_GREEN))

# 右卡 Neo4j
card(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(4.0), C_BG3, C_PURPLE)
rect(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.7), C_BG2)
icon(s, Inches(7.3), Inches(1.95), Inches(0.5), "🔗", C_BG2, C_PURPLE)
text(s, Inches(8.0), Inches(1.95), Inches(4.5), Inches(0.5), "路径2: Neo4j 知识图谱", size=16, bold=True, color=C_PURPLE)
items_neo = [
    ("按字段精确匹配 DTC/车型/场景", C_PURPLE),
    ("2036 节点 / 4356 关系", C_LIGHT),
    ("字段分通道独立匹配 零拼接", C_LIGHT),
    ("适配: 有确定字段时零误召", C_GREEN)
]
for i, (item, c) in enumerate(items_neo):
    y = Inches(2.75 + i * 0.6)
    rect(s, Inches(7.4), y + Inches(0.08), Inches(0.08), Inches(0.08), c)
    text(s, Inches(7.6), y, Inches(4.8), Inches(0.5), item, size=14, color=c, bold=(c==C_GREEN))

# 底部融合箭头
card(s, Inches(2.5), Inches(6.1), Inches(8.3), Inches(0.85), C_BG2, C_GREEN)
text(s, Inches(2.5), Inches(6.1), Inches(8.3), Inches(0.85),
     "→ 结果合并去重 → Embedding 精排 → 传给 LLM ←", size=16, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ========== 第5页 六层架构 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "03", "项目搭建思路：LangChain 架构", C_BLUE)

layers = [
    ("适配层", "adapter/", "FastAPI 对接上游平台（异步提交+回调）", C_CYAN),
    ("Agent 层", "agent/", "ReAct 推理循环，编排工具调用", C_BLUE),
    ("检索层", "retrieval/", "HybridRetriever 双路召回 + 精排", C_PURPLE),
    ("知识层", "knowledge/", "对话抽取实体关系 → 审核 → 写回图谱", C_GREEN),
    ("上下文层", "context/", "三层记忆 + 话题检测 + 异步摘要", C_ORANGE),
    ("存储层", "storage/", "Neo4j + ChromaDB + Redis + 磁盘", C_PINK)
]
for i, (name, path, desc, color) in enumerate(layers):
    y = Inches(1.4 + i * 0.88)
    # 主卡片
    card(s, Inches(0.5), y, Inches(12.3), Inches(0.75), C_BG3)
    # 左侧色条
    rect(s, Inches(0.5), y, Inches(0.08), Inches(0.75), color)
    # 层名
    rect(s, Inches(0.8), y + Inches(0.12), Inches(1.6), Inches(0.5), C_BG2)
    text(s, Inches(0.8), y + Inches(0.12), Inches(1.6), Inches(0.5), name, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    # 路径
    text(s, Inches(2.7), y + Inches(0.15), Inches(2.5), Inches(0.5), path, size=13, color=C_CYAN, font="Consolas")
    # 描述
    text(s, Inches(5.5), y + Inches(0.15), Inches(7), Inches(0.5), desc, size=14, color=C_LIGHT)
    # 右侧色点
    icon(s, Inches(12.2), y + Inches(0.15), Inches(0.4), "", color, C_WHITE)

# ========== 第6页 Agent 工具集 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "03", "Agent 工具集（动态注册）", C_BLUE)
text(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
     "通过 LangChain StructuredTool 暴露 3 个工具，不可用时自动降级不注册", size=14, color=C_MUTE)

tools = [
    ("search_similar_incidents", "模糊匹配故障现象", "ChromaDB 向量检索", "🔍", C_CYAN),
    ("query_fault_graph", "结构化字段精确查", "Neo4j Cypher 查询", "🔗", C_PURPLE),
    ("get_incident_detail", "查看工单详情", "历史记录", "📄", C_GREEN)
]
for i, (name, trig, src, icon_sym, color) in enumerate(tools):
    y = Inches(2.0 + i * 1.4)
    card(s, Inches(0.5), y, Inches(12.3), Inches(1.2), C_BG3)
    rect(s, Inches(0.5), y, Inches(0.08), Inches(1.2), color)
    # 图标
    icon(s, Inches(0.9), y + Inches(0.3), Inches(0.6), icon_sym, C_BG2, color)
    # 工具名
    text(s, Inches(1.8), y + Inches(0.2), Inches(4.5), Inches(0.5), name, size=15, bold=True, color=color, font="Consolas")
    # 触发场景
    text(s, Inches(1.8), y + Inches(0.7), Inches(4.5), Inches(0.4), f"触发: {trig}", size=12, color=C_LIGHT)
    # 数据源 chip
    rounded(s, Inches(6.8), y + Inches(0.35), Inches(3.5), Inches(0.4), C_BG2)
    text(s, Inches(6.8), y + Inches(0.35), Inches(3.5), Inches(0.4), f"数据源: {src}", size=12, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 右侧装饰
    text(s, Inches(11), y + Inches(0.4), Inches(1.5), Inches(0.5), "✓ 已注册", size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ========== 第7页 端到端流程 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "04", "实际工作流：端到端流程", C_PURPLE)

steps = [
    ("1", "用户输入", "MCU报P1A3E98爬坡IGBT过温", "📝", C_CYAN),
    ("2", "话题命中判别", "排除非诊断话题 关联上下文", "🧠", C_BLUE),
    ("3", "双路预检索", "Chroma 5条 + Neo4j 3条 精排", "🔍", C_PURPLE),
    ("4", "ReAct 推理循环", "Thought → Action → Observation → Final", "⚙️", C_ORANGE),
    ("5", "可选双层输出", "Markdown 人读 + CSV/JSON 机读", "📊", C_GREEN),
    ("6", "知识沉淀", "LLM 抽取实体/关系 → 待审队列", "💾", C_PINK)
]
for i, (n, t, d, icon_sym, color) in enumerate(steps):
    y = Inches(1.35 + i * 0.92)
    # 卡片
    card(s, Inches(0.5), y, Inches(12.3), Inches(0.8), C_BG3)
    rect(s, Inches(0.5), y, Inches(0.08), Inches(0.8), color)
    # 序号圆
    icon(s, Inches(0.8), y + Inches(0.1), Inches(0.6), n, color, C_WHITE)
    # 图标
    text(s, Inches(1.6), y + Inches(0.15), Inches(0.5), Inches(0.5), icon_sym, size=20, align=PP_ALIGN.CENTER)
    # 标题
    text(s, Inches(2.3), y + Inches(0.1), Inches(3.5), Inches(0.4), t, size=15, bold=True, color=C_WHITE)
    # 描述
    text(s, Inches(6), y + Inches(0.15), Inches(6.5), Inches(0.4), d, size=13, color=C_LIGHT)
    # 箭头（除最后一个）
    if i < len(steps) - 1:
        text(s, Inches(0.8), y + Inches(0.75), Inches(0.6), Inches(0.2), "↓", size=14, color=C_MUTE, align=PP_ALIGN.CENTER)

# ========== 第8页 CAN 兜底 + 三层记忆 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "04", "CAN 报文兜底 + 三层记忆系统", C_ORANGE)

# 左：CAN 兜底
card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), C_BG3, C_ORANGE)
text(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5), "CAN 报文自动兜底", size=18, bold=True, color=C_ORANGE)
icon(s, Inches(0.8), Inches(2.1), Inches(0.6), "📡", C_BG2, C_ORANGE)
text(s, Inches(1.6), Inches(2.15), Inches(4), Inches(0.5), "当预检索结果不足（top1 < 0.6）", size=13, color=C_LIGHT)

can_steps = [
    ("预检索相似度低", C_PINK),
    ("解码 CAN 工况文件", C_CYAN),
    ("支持 ASC/BLF/MF4/CSV", C_MUTE),
    ("DBC 解码为物理值", C_CYAN),
    ("信号摘要注入上下文", C_GREEN)
]
for i, (step, color) in enumerate(can_steps):
    y = Inches(2.9 + i * 0.6)
    rect(s, Inches(1.0), y + Inches(0.1), Inches(0.12), Inches(0.12), color)
    text(s, Inches(1.3), y, Inches(4.5), Inches(0.5), step, size=14, color=C_LIGHT if color == C_MUTE else color)
    if i < len(can_steps) - 1:
        text(s, Inches(1.0), y + Inches(0.4), Inches(0.3), Inches(0.2), "↓", size=12, color=C_MUTE)

# 右：三层记忆
card(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), C_BG3, C_BLUE)
text(s, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5), "三层记忆系统", size=18, bold=True, color=C_BLUE)

mem_layers = [
    ("热层", "完整消息", "Redis", "🔥", C_PINK),
    ("温层", "滚动摘要", "Redis", "🌡️", C_ORANGE),
    ("冷层", "归档持久化", "磁盘", "❄️", C_CYAN)
]
for i, (name, content, storage, icon_sym, color) in enumerate(mem_layers):
    y = Inches(2.3 + i * 1.4)
    # 子卡片
    rounded(s, Inches(7.3), y, Inches(5.2), Inches(1.15), C_BG2)
    rect(s, Inches(7.3), y, Inches(0.08), Inches(1.15), color)
    # 图标
    text(s, Inches(7.6), y + Inches(0.3), Inches(0.5), Inches(0.5), icon_sym, size=22, align=PP_ALIGN.CENTER)
    # 层名
    text(s, Inches(8.3), y + Inches(0.15), Inches(1.5), Inches(0.5), name, size=16, bold=True, color=color)
    # 内容
    text(s, Inches(8.3), y + Inches(0.6), Inches(2.5), Inches(0.4), content, size=12, color=C_LIGHT)
    # 存储 chip
    rounded(s, Inches(11.0), y + Inches(0.35), Inches(1.3), Inches(0.35), C_BG3)
    text(s, Inches(11.0), y + Inches(0.35), Inches(1.3), Inches(0.35), storage, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

text(s, Inches(7.3), Inches(6.5), Inches(5.2), Inches(0.3),
     "Redis 不可用时自动降级内存，跨会话可从冷层恢复", size=11, color=C_MUTE, align=PP_ALIGN.CENTER)

# ========== 第9页 六大创新点 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "05", "核心创新点", C_GREEN)
text(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
     "与普通 RAG 相比的六大创新", size=14, color=C_MUTE)

innovs = [
    ("检索模式", "单路向量", "双路混合 RAG", C_CYAN),
    ("LLM 角色", "上下文生成器", "ReAct 推理编排者", C_BLUE),
    ("知识库", "静态", "动态闭环（自动抽取→审核→写入）", C_PURPLE),
    ("多轮记忆", "无/简单", "三层分层记忆（热/温/冷）", C_GREEN),
    ("兜底机制", "无", "CAN 报文自动解码注入信号", C_ORANGE),
    ("输出格式", "纯文本", "双层输出（Markdown + JSON）", C_PINK)
]
# 表头
card(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5), C_BG2)
text(s, Inches(0.8), Inches(1.72), Inches(2.5), Inches(0.45), "创新点", size=13, bold=True, color=C_WHITE)
text(s, Inches(4.0), Inches(1.72), Inches(3.5), Inches(0.45), "普通 RAG", size=13, bold=True, color=C_MUTE)
text(s, Inches(8.0), Inches(1.72), Inches(4.5), Inches(0.45), "我们的方案", size=13, bold=True, color=C_GREEN)

for i, (a, b, c, color) in enumerate(innovs):
    y = Inches(2.3 + i * 0.7)
    bg = C_BG3 if i % 2 == 0 else C_BG2
    rounded(s, Inches(0.5), y, Inches(12.3), Inches(0.6), bg)
    rect(s, Inches(0.5), y, Inches(0.06), Inches(0.6), color)
    text(s, Inches(0.8), y + Inches(0.1), Inches(3.0), Inches(0.4), a, size=14, bold=True, color=C_WHITE)
    text(s, Inches(4.0), y + Inches(0.1), Inches(3.5), Inches(0.4), b, size=13, color=C_MUTE)
    text(s, Inches(8.0), y + Inches(0.1), Inches(4.5), Inches(0.4), c, size=13, bold=True, color=color)

# ========== 第10页 效果展示 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "06", "效果展示：一次完整诊断输出", C_PINK)

# 左：Markdown 报告
card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), C_BG3, C_GREEN)
rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.6), C_BG2)
icon(s, Inches(0.7), Inches(1.35), Inches(0.5), "📄", C_BG2, C_GREEN)
text(s, Inches(1.4), Inches(1.35), Inches(4), Inches(0.5), "人类可读 Markdown 报告", size=14, bold=True, color=C_GREEN)

md_items = [
    ("故障分类", "过温故障", C_CYAN),
    ("根因分析", "IGBT 散热基板焊接不良", C_PINK),
    ("", "动态行驶螺丝松动与PCB短路", C_MUTE),
    ("解决方案", "优化作业流程", C_GREEN),
    ("", "螺纹孔深度纳入作业管理", C_MUTE),
    ("相似工单", "5 条（车型/工况/软件版本）", C_ORANGE),
    ("仪表指示灯", "电机故障红灯", C_PURPLE)
]
y_off = 2.1
for label, val, color in md_items:
    if label:
        text(s, Inches(0.9), y_off, Inches(1.5), Inches(0.4), f"• {label}", size=12, bold=True, color=color)
        text(s, Inches(2.5), y_off, Inches(3.5), Inches(0.4), val, size=12, color=C_WHITE)
    else:
        text(s, Inches(2.5), y_off, Inches(3.5), Inches(0.4), val, size=11, color=color)
    y_off += 0.55

# 右：JSON 代码块
card(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), RGBColor(0x0a, 0x0f, 0x1a), C_PINK)
rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.6), C_BG2)
icon(s, Inches(7.2), Inches(1.35), Inches(0.5), "{ }", C_BG2, C_PINK)
text(s, Inches(7.9), Inches(1.35), Inches(4), Inches(0.5), "机器可读 JSON 条目", size=14, bold=True, color=C_PINK)

code_lines = [
    ('{', C_MUTE),
    ('  "classification":', C_CYAN), ('"过温故障",', C_GREEN),
    ('  "root_cause":', C_CYAN), ('"IGBT散热基板焊接不良",', C_GREEN),
    ('  "solution":', C_CYAN), ('"优化作业流程",', C_GREEN),
    ('  "dtc_code":', C_CYAN), ('"P1A3E98",', C_GREEN),
    ('  "similar_record_ids":', C_CYAN), ['[...]', C_ORANGE],
    ('  "confidence":', C_CYAN), ('0.91', C_ORANGE),
    ('}', C_MUTE)
]
y_off = 2.2
for line in code_lines:
    if isinstance(line, list):
        txt, c = line[0], line[1]
    else:
        txt, c = line
    text(s, Inches(7.3), y_off, Inches(5.2), Inches(0.4), txt, size=13, color=c, font="Consolas")
    y_off += 0.35

# 底部说明
text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
     "故障分类强制从 10 类固定列表选择，保证分类一致性", size=12, color=C_MUTE, align=PP_ALIGN.CENTER)

# ========== 第11页 数据规模 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "06", "数据规模与测试", C_PINK)

cards_data = [
    ("87", "测试用例", "全部通过 41 秒", C_GREEN, "✓"),
    ("628", "向量库记录", "1024 维 text-embedding-v4", C_CYAN, "📊"),
    ("2036", "知识图谱节点", "4356 关系 / 8 类实体", C_PURPLE, "🕸"),
    ("8000", "Token 预算", "热层窗口自适应 2~20 轮", C_ORANGE, "⚡")
]
for i, (num, label, sub, color, icon_sym) in enumerate(cards_data):
    x = Inches(0.4 + i * 3.25)
    card(s, x, Inches(1.5), Inches(3.0), Inches(3.5), C_BG3, color)
    # 顶部图标
    icon(s, x + Inches(1.0), Inches(1.7), Inches(0.9), icon_sym, C_BG2, color)
    # 大数字
    text(s, x, Inches(2.7), Inches(3.0), Inches(1.0), num, size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 标签
    text(s, x, Inches(3.7), Inches(3.0), Inches(0.4), label, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    # 描述
    text(s, x + Inches(0.2), Inches(4.1), Inches(2.6), Inches(0.8), sub, size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

# 性能指标
text(s, Inches(0.5), Inches(5.3), Inches(12), Inches(0.4), "性能指标", size=16, bold=True, color=C_CYAN)
perfs = [
    ("LLM 推理响应", "0.4~1.0 秒", "单轮 e2e 实测", C_CYAN),
    ("检索 top1 相似度", "0.71~0.93", "符合预期区间", C_GREEN),
    ("温层摘要压缩", "Token 减少 60%", "合并后显著下降", C_ORANGE)
]
for i, (k, v, d, color) in enumerate(perfs):
    x = Inches(0.4 + i * 4.2)
    rounded(s, x, Inches(5.8), Inches(3.9), Inches(1.0), C_BG3)
    rect(s, x, Inches(5.8), Inches(0.06), Inches(1.0), color)
    text(s, x + Inches(0.3), Inches(5.9), Inches(3.5), Inches(0.35), k, size=12, color=C_MUTE)
    text(s, x + Inches(0.3), Inches(6.2), Inches(3.5), Inches(0.4), v, size=18, bold=True, color=color)
    text(s, x + Inches(0.3), Inches(6.55), Inches(3.5), Inches(0.35), d, size=11, color=C_LIGHT)

# ========== 第12页 技术栈 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "07", "技术栈", C_ORANGE)

techs = [
    ("LLM 框架", "LangChain create_agent ReAct", "🧩", C_CYAN),
    ("LLM 后端", "阿里云 DashScope 通义千问", "☁️", C_BLUE),
    ("知识图谱", "Neo4j 5.22.0 Cypher + APOC", "🕸", C_PURPLE),
    ("向量数据库", "ChromaDB cosine PersistentClient", "📐", C_GREEN),
    ("会话存储", "Redis 热温层 + 磁盘冷层归档", "💾", C_ORANGE),
    ("Web 框架", "FastAPI + Uvicorn", "🚀", C_PINK),
    ("CLI", "Typer + Rich 交互式终端", "⌨️", C_CYAN),
    ("语言", "Python 3.10+", "🐍", C_GREEN)
]
for i, (cat, name, icon_sym, color) in enumerate(techs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.4 + row * 1.4)
    card(s, x, y, Inches(6.0), Inches(1.2), C_BG3)
    rect(s, x, y, Inches(0.08), Inches(1.2), color)
    # 图标
    icon(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), icon_sym, C_BG2, color)
    # 分类
    text(s, x + Inches(1.1), y + Inches(0.15), Inches(4.5), Inches(0.4), cat, size=12, color=C_MUTE)
    # 名称
    text(s, x + Inches(1.1), y + Inches(0.55), Inches(4.5), Inches(0.5), name, size=15, bold=True, color=C_WHITE)

# ========== 第13页 演示命令 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
header(s, "07", "演示命令", C_ORANGE)

cmds = [
    ("交互式多轮诊断", "python -m diagnosis_agent.cli chat", C_CYAN),
    ("单次文本诊断", "python -m diagnosis_agent.cli diagnose --text ...", C_GREEN),
    ("加载 CSV 工单数据", "python -m diagnosis_agent.cli load-data --file ...", C_ORANGE),
    ("标准接口对接平台", "python -m diagnosis_agent.cli diagnose --json-input ...", C_PURPLE),
    ("查看知识库统计", "python -m diagnosis_agent.cli stats", C_PINK),
    ("启动 FastAPI 服务", "python -m diagnosis_agent.cli adapter", C_BLUE)
]
for i, (desc, cmd, color) in enumerate(cmds):
    y = Inches(1.4 + i * 0.9)
    # 终端风格卡片
    card(s, Inches(0.5), y, Inches(12.3), Inches(0.75), RGBColor(0x0a, 0x0f, 0x1a))
    rect(s, Inches(0.5), y, Inches(0.08), Inches(0.75), color)
    # 终端圆点
    for dx, dc in [(Inches(0.9), RGBColor(0xff,0x5f,0x57)), (Inches(1.1), RGBColor(0xff,0xbd,0x2d)), (Inches(1.3), RGBColor(0x23,0xc8,0x3b))]:
        rect(s, dx, y + Inches(0.25), Inches(0.15), Inches(0.15), dc, shape=MSO_SHAPE.OVAL)
    # 注释
    text(s, Inches(1.7), y + Inches(0.2), Inches(3.5), Inches(0.4), f"# {desc}", size=12, color=C_MUTE)
    # 命令
    text(s, Inches(5.5), y + Inches(0.2), Inches(7), Inches(0.4), cmd, size=13, color=color, font="Consolas")

# ========== 第14页 结束页 ==========
s = prs.slides.add_slide(BLANK)
slide_dark(s)
# 背景渐变色块
rect(s, 0, 0, SW, SH, C_BG1)
rect(s, 0, Inches(5.0), SW, Inches(2.5), C_BG2)
# 顶部装饰
rect(s, 0, Inches(3.0), SW, Inches(0.06), C_CYAN)
# 大标题
text(s, 0, Inches(2.0), SW, Inches(1.2), "谢谢", size=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
text(s, 0, Inches(3.2), SW, Inches(0.8), "欢迎各位老师提问", size=26, color=C_CYAN, align=PP_ALIGN.CENTER)
# 底部小字
text(s, 0, Inches(5.5), SW, Inches(0.4), "Diagnosis Agent v0.5.0", size=14, color=C_MUTE, align=PP_ALIGN.CENTER)
text(s, 0, Inches(6.0), SW, Inches(0.4), "LLM + 双路 RAG + 知识图谱闭环", size=13, color=C_MUTE, align=PP_ALIGN.CENTER)

out_path = "/home/dfmc/diagnose_agent/shared/diagnosis_agent/docs/答辩文档_v3.pptx"
prs.save(out_path)
print(f"已生成: {out_path}")
print(f"共 {len(prs.slides)} 页")
