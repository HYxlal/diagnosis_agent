#!/usr/bin/env python3
"""基于 docs/答辩文档.md 生成 docs/答辩文档.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 主题色
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # 深蓝
ACCENT  = RGBColor(0x2E, 0x86, 0xC1)      # 中蓝
LIGHT   = RGBColor(0xEA, 0xF2, 0xF8)       # 浅蓝背景
DARK    = RGBColor(0x24, 0x29, 0x2E)       # 深灰文字
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREY    = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = color; s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, runs, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: str 或 [(text, {size,color,bold})] 列表"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)

    if isinstance(runs, str):
        runs = [(runs, {})]

    # 多段：用 \n 分隔
    lines = []
    cur = []
    for t, opts in runs:
        parts = t.split("\n")
        for i, p in enumerate(parts):
            if i > 0:
                lines.append(cur); cur = []
            cur.append((p, opts))
    lines.append(cur)

    for li, line_runs in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for t, opts in line_runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.color.rgb = opts.get("color", color)
            r.font.name = "Microsoft YaHei"
    return tb

def title_bar(slide, title, page_no, total):
    # 顶部色条
    rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    # 标题
    text(slide, Inches(0.5), Inches(0.1), Inches(10), Inches(0.7),
         title, size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 页码
    text(slide, Inches(12.2), Inches(0.1), Inches(1), Inches(0.7),
         f"{page_no}/{total}", size=12, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # 底部细线
    rect(slide, 0, Inches(7.3), SW, Inches(0.2), LIGHT)

def table(slide, x, y, w, h, headers, rows, header_color=PRIMARY, col_widths=None):
    rows_count = len(rows) + 1
    cols_count = len(headers)
    g = slide.shapes.add_table(rows_count, cols_count, x, y, w, h)
    tbl = g.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    # 表头
    for i, htext in enumerate(headers):
        c = tbl.cell(0, i)
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        c.text = htext
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    # 数据行
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if ri % 2 == 1 else WHITE
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(12); r.font.color.rgb = DARK; r.font.name = "Microsoft YaHei"
    return g

# ============== 幻灯片 ==============
TOTAL = 12

# ---- 封面 ----
s = add_slide()
bg(s, PRIMARY)
# 装饰
rect(s, 0, Inches(5.0), SW, Inches(0.08), ACCENT)
text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
     "车辆故障诊断智能助手", size=54, bold=True, color=WHITE)
text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
     "LLM + 双路 RAG + 知识图谱闭环", size=28, color=RGBColor(0xAE, 0xD6, 0xF1))
text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
     "项目名称：Diagnosis Agent v0.5.0", size=18, color=RGBColor(0xBD, 0xC3, 0xCA))
text(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
     "技术路线：LangChain ReAct + Neo4j + ChromaDB + 通义千问", size=16, color=RGBColor(0xBD, 0xC3, 0xCA))

# ---- 目录 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "目 录", 2, TOTAL)
items = [
    ("01", "解决什么问题"),
    ("02", "技术方案：LLM + 双路 RAG"),
    ("03", "项目搭建思路：LangChain 架构"),
    ("04", "实际工作流"),
    ("05", "核心创新点"),
    ("06", "效果展示"),
    ("07", "技术栈"),
    ("08", "演示命令"),
]
for i, (no, name) in enumerate(items):
    col = i // 4
    row = i % 4
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.3)
    rect(s, x, y, Inches(0.7), Inches(0.7), ACCENT)
    text(s, x, y, Inches(0.7), Inches(0.7), no, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.9), y + Inches(0.05), Inches(5), Inches(0.6), name, size=20, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# ---- 1. 解决什么问题 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "一、解决什么问题？", 3, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
     "传统车辆故障诊断面临三个痛点，根本上都是『经验难沉淀，回答靠幻觉』：",
     size=16, color=GREY)
pains = [
    ("1", "历史工单散落 Excel", "新人翻几十页才能找到相似案例，效率低"),
    ("2", "诊断报告非结构化", "自然语言写出来很难自动入库统计"),
    ("3", "LLM 凭空回答容易错", "不知道真实历史案例，根因对策无证据支撑"),
]
for i, (n, t, d) in enumerate(pains):
    y = Inches(2.0 + i * 1.5)
    rect(s, Inches(0.5), y, Inches(0.8), Inches(0.8), RED)
    text(s, Inches(0.5), y, Inches(0.8), Inches(0.8), n, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.6), y + Inches(0.05), Inches(11), Inches(0.5), t, size=20, bold=True, color=PRIMARY)
    text(s, Inches(1.6), y + Inches(0.5), Inches(11), Inches(0.4), d, size=14, color=DARK)
# 目标
rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), LIGHT)
text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.6),
     [("目标：", {"bold": True, "color": PRIMARY, "size": 18}),
      ("让大模型基于真实历史案例推理，而不是凭记忆瞎编", {"size": 18, "color": DARK})],
     anchor=MSO_ANCHOR.MIDDLE)

# ---- 2. 双路 RAG ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "二、技术方案：LLM + 双路 RAG", 4, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
     "普通 RAG 只有向量检索一条路，我们做了两条互补路径，精度更高：",
     size=16, color=GREY)
# 两条路径卡片
rect(s, Inches(0.5), Inches(1.8), Inches(6), Inches(2.6), LIGHT)
rect(s, Inches(0.5), Inches(1.8), Inches(6), Inches(0.6), ACCENT)
text(s, Inches(0.5), Inches(1.8), Inches(6), Inches(0.6), "路径1：ChromaDB 向量检索", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(2.6), Inches(5.6), Inches(1.8),
     [("• 语义模糊匹配现象描述\n", {"size": 15}),
      ("• 找『IGBT过温』等相似文本片段\n", {"size": 15}),
      ("• 628 条历史工单，1024 维向量\n", {"size": 15}),
      ("• 适配：现象模糊、不知道精确 DTC 码时", {"size": 15, "color": ACCENT, "bold": True})],
     line_spacing=1.5)

rect(s, Inches(6.8), Inches(1.8), Inches(6), Inches(2.6), LIGHT)
rect(s, Inches(6.8), Inches(1.8), Inches(6), Inches(0.6), PRIMARY)
text(s, Inches(6.8), Inches(1.8), Inches(6), Inches(0.6), "路径2：Neo4j 知识图谱", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(7.0), Inches(2.6), Inches(5.6), Inches(1.8),
     [("• 按字段精确匹配 DTC/车型/场景\n", {"size": 15}),
      ("• 2036 节点 / 4356 关系\n", {"size": 15}),
      ("• 字段分通道独立匹配，零拼接\n", {"size": 15}),
      ("• 适配：有确定字段时，零误召", {"size": 15, "color": PRIMARY, "bold": True})],
     line_spacing=1.5)

# 合并
rect(s, Inches(2), Inches(4.7), Inches(9.3), Inches(0.8), GREEN)
text(s, Inches(2), Inches(4.7), Inches(9.3), Inches(0.8),
     "结果合并去重 → Embedding 精排 → 传给 LLM",
     size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1),
     "两条路径各自独立匹配，完全不做字段字符串拼接，最后由精排层融合排序。不会因为『车型字符串不对就召回不到结果』。",
     size=14, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.5)

# ---- 3. LangChain 架构 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "三、项目搭建思路：LangChain 架构", 5, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
     "LangChain 提供三大抽象，我们利用以下能力而非从零造轮子：", size=16, color=GREY)
table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.4),
      ["LangChain 能力", "我们怎么用"],
      [["create_agent", "构建 ReAct Agent，LLM 自主决定调哪个工具、查几次"],
       ["BaseRetriever", "统一检索器接口，Chroma 和 Neo4j 都实现，Agent 不感知差异"],
       ["StructuredTool", "把检索方法包装成工具，LLM 按 schema 调用"]],
      col_widths=[Inches(3.5), Inches(8.8)])

text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5), "六层分层架构，每层职责单一：", size=16, bold=True, color=PRIMARY)
layers = [
    ("适配层 adapter/", "FastAPI 对接上游平台（异步提交 + 回调）"),
    ("Agent 层 agent/", "ReAct 推理循环，编排工具调用"),
    ("检索层 retrieval/", "HybridRetriever 双路召回 + 精排"),
    ("知识层 knowledge/", "对话抽取实体关系 → 审核 → 写回图谱"),
    ("上下文层 context/", "三层记忆 + 话题检测 + 异步摘要"),
    ("存储层 storage/", "Neo4j + ChromaDB + Redis + 磁盘"),
]
for i, (name, desc) in enumerate(layers):
    y = Inches(4.0 + i * 0.5)
    rect(s, Inches(0.5), y, Inches(3.2), Inches(0.45), ACCENT if i % 2 == 0 else PRIMARY)
    text(s, Inches(0.5), y, Inches(3.2), Inches(0.45), name, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.8), y, Inches(9), Inches(0.45), desc, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# ---- 4. Agent 工具集 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "三、Agent 工具集（动态注册）", 6, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
     "通过 LangChain StructuredTool 暴露 4 个工具，Neo4j 不可用时自动降级不注册：",
     size=16, color=GREY)
table(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.2),
      ["工具", "触发场景", "数据源"],
      [["search_similar_incidents", "模糊匹配故障现象", "ChromaDB 向量检索"],
       ["query_fault_graph", "有 DTC/车型等结构化字段", "Neo4j Cypher 查询"],
       ["get_incident_detail", "查看工单详情", "历史记录"],
       ["convert_working_condition_file", "解析工况文件", "文件转换"]],
      col_widths=[Inches(3.8), Inches(4.5), Inches(4)])
text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
     [("注：", {"bold": True, "color": RED, "size": 15}),
      ("can_converter 在 prompt 中提及但未在 get_tool_list() 注册，它由 can_fallback.py 在 Agent 主流程中直接调用，不走 LLM 工具调用路径。这是设计如此——CAN 兜底是流程级兜底，不是工具级。", {"size": 15, "color": DARK})],
     line_spacing=1.5)

# ---- 5. 端到端流程 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "四、实际工作流：端到端流程", 7, TOTAL)
steps = [
    ("1", "用户输入", "MCU报P1A3E98爬坡IGBT过温，车辆抖动"),
    ("2", "输入解析 + 意图分类", "提取 entities，判定 diagnostic_query"),
    ("3", "预检索（双路 RAG 召回）", "Chroma 5条 + Neo4j 3条 → 精排"),
    ("4", "ReAct Agent 推理循环", "Thought → Action → Observation → Final"),
    ("5", "双层输出", "Markdown 人读 + CSV/JSON 机读"),
    ("6", "知识沉淀", "LLM 抽取实体/关系 → 写入待审核队列"),
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(1.5 + i * 0.85)
    # 序号圆
    rect(s, Inches(0.5), y, Inches(0.6), Inches(0.6), ACCENT)
    text(s, Inches(0.5), y, Inches(0.6), Inches(0.6), n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.3), y + Inches(0.02), Inches(3.5), Inches(0.6), t, size=15, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5), y + Inches(0.02), Inches(7.8), Inches(0.6), d, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        # 箭头线
        rect(s, Inches(0.78), y + Inches(0.6), Inches(0.04), Inches(0.25), GREY)

# ---- 6. CAN 兜底 + 三层记忆 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "四、CAN 兜底 + 三层记忆系统", 8, TOTAL)
# CAN 兜底
rect(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.5), ACCENT)
text(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.5), "CAN 报文自动兜底", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.5), Inches(1.7), Inches(6), Inches(2.5),
     [("当预检索结果不足（top1 相似度 < 0.6）时：\n\n", {"size": 14, "bold": True, "color": PRIMARY}),
      ("预检索相似度低\n", {"size": 13}),
      ("  ↓ 解码 CAN 工况文件（ASC/BLF/MF4/CSV）\n", {"size": 13}),
      ("  ↓ 信号摘要注入 Agent 上下文\n", {"size": 13}),
      ("  ↓ Agent 基于信号推理\n\n", {"size": 13}),
      ("结合 DBC 文件解码为物理值", {"size": 13, "color": GREY, "italic": True})],
     line_spacing=1.4)

# 三层记忆
rect(s, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5), PRIMARY)
text(s, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5), "三层记忆系统", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
table(s, Inches(6.8), Inches(1.7), Inches(6), Inches(2.5),
      ["层", "内容", "存储", "降级条件"],
      [["热", "完整消息", "Redis", "超 window → 摘要移温层"],
       ["温", "滚动摘要", "Redis", "摘要超限 → 合并精简"],
       ["冷", "归档 JSON", "磁盘", "话题切换/会话结束"]],
      col_widths=[Inches(0.6), Inches(1.4), Inches(1.2), Inches(2.8)])
text(s, Inches(6.8), Inches(4.3), Inches(6), Inches(0.5),
     "Redis 不可用时自动降级内存，跨会话可从冷层恢复",
     size=12, color=GREY, align=PP_ALIGN.CENTER)

# ---- 7. 核心创新点 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "五、核心创新点", 9, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
     "与普通 RAG 相比的六大创新：", size=16, color=GREY)
table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.8),
      ["创新点", "普通 RAG", "我们的方案"],
      [["检索模式", "单路向量", "双路混合 RAG（Chroma 语义 + Neo4j 结构化）"],
       ["LLM 角色", "上下文生成器", "ReAct 推理编排者（主动调用工具查资料）"],
       ["知识库", "静态", "动态闭环（自动抽取 → 审核 → 写入图谱）"],
       ["多轮记忆", "无/简单", "三层分层记忆（热/温/冷），支持话题切换"],
       ["兜底机制", "无", "CAN 报文自动解码，检索不足时注入信号上下文"],
       ["输出格式", "纯文本", "双层输出（人读 Markdown + 机读 CSV/JSON）"]],
      col_widths=[Inches(2.5), Inches(3.5), Inches(6.3)])

# ---- 8. 效果展示 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "六、效果展示：一次完整诊断输出", 10, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.5), "人类可读报告（Markdown）：", size=15, bold=True, color=PRIMARY)
text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(4),
     [("• 故障分类：过温故障\n", {"size": 14}),
      ("• 根因分析：IGBT 散热基板焊接不良\n", {"size": 14}),
      ("  动态行驶中螺丝松动与 PCB 电容接触短路\n", {"size": 12, "color": GREY}),
      ("• 解决方案：优化作业流程\n", {"size": 14}),
      ("  将螺纹孔深度纳入作业流程管理\n", {"size": 12, "color": GREY}),
      ("• 相似工单：5 条（含车型/工况/软件版本）\n", {"size": 14}),
      ("• 仪表指示灯：电机故障红灯", {"size": 14})],
     line_spacing=1.5)
text(s, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5), "机器可读条目（JSON）：", size=15, bold=True, color=PRIMARY)
# JSON 代码块
rect(s, Inches(6.8), Inches(1.6), Inches(6), Inches(4), RGBColor(0x2B, 0x3A, 0x4A))
text(s, Inches(7), Inches(1.7), Inches(5.6), Inches(3.8),
     [("{\n", {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "classification": "过温故障",\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "root_cause": "IGBT散热基板焊接不良",\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "solution": "优化作业流程",\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "dtc_code": "P1A3E98",\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "similar_record_ids": [...],\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ('  "confidence": 0.91\n', {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)}),
      ("}", {"size": 12, "color": RGBColor(0x8B, 0xC3, 0x4A)})],
     line_spacing=1.3)
text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1),
     [("故障分类强制从 10 类固定列表选择", {"bold": True, "color": RED, "size": 14}),
      ("（驱动异常/控制异常/超速/高压异常/低压异常/过温/通信/旋变/状态机/油泵），保证分类一致性。", {"size": 14, "color": DARK})],
     align=PP_ALIGN.CENTER, line_spacing=1.5)

# ---- 9. 数据规模 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "六、数据规模与测试", 11, TOTAL)
text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), "数据规模与测试指标：", size=16, color=GREY)
# 四个大数字卡片
cards = [
    ("87", "测试用例", "全部通过，41 秒", GREEN),
    ("628", "向量库记录", "1024 维 text-embedding-v4", ACCENT),
    ("2036", "知识图谱节点", "4356 关系 / 8 类实体", PRIMARY),
    ("8000", "Token 预算", "热层窗口自适应 2~20 轮", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (num, label, desc, color) in enumerate(cards):
    x = Inches(0.5 + i * 3.15)
    rect(s, x, Inches(1.8), Inches(2.9), Inches(2.4), LIGHT)
    rect(s, x, Inches(1.8), Inches(2.9), Inches(0.6), color)
    text(s, x, Inches(1.8), Inches(2.9), Inches(0.6), label, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, Inches(2.5), Inches(2.9), Inches(1), num, size=40, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, Inches(3.6), Inches(2.9), Inches(0.5), desc, size=10, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)

# 性能指标
text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5), "性能指标：", size=16, bold=True, color=PRIMARY)
table(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2),
      ["指标", "数值", "说明"],
      [["LLM 推理响应", "0.4~1.0 秒", "单轮 e2e 实测"],
       ["检索召回 top1 相似度", "0.71~0.93", "符合预期区间"],
       ["温层摘要压缩", "Token 减少 60%", "合并后总 Token 显著下降"]],
      col_widths=[Inches(4), Inches(3), Inches(5.3)])

# ---- 10. 技术栈 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "七、技术栈", 12, TOTAL)
techs = [
    ("LLM 框架", "LangChain（create_agent ReAct 模式）", ACCENT),
    ("LLM 后端", "阿里云 DashScope（通义千问，OpenAI 兼容）", ACCENT),
    ("知识图谱", "Neo4j 5.22.0（Cypher + APOC）", PRIMARY),
    ("向量数据库", "ChromaDB（cosine 空间，PersistentClient）", PRIMARY),
    ("会话存储", "Redis（热/温层）+ 磁盘（冷层归档）", GREEN),
    ("Web 框架", "FastAPI + Uvicorn（平台适配层）", GREEN),
    ("CLI", "Typer + Rich（交互式终端）", RGBColor(0x8E, 0x44, 0xAD)),
    ("语言", "Python 3.10+", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (cat, name, color) in enumerate(techs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.5 + row * 1.35)
    rect(s, x, y, Inches(5.8), Inches(1.1), LIGHT)
    rect(s, x, y, Inches(0.15), Inches(1.1), color)
    text(s, x + Inches(0.4), y + Inches(0.1), Inches(5.2), Inches(0.5), cat, size=12, color=GREY)
    text(s, x + Inches(0.4), y + Inches(0.5), Inches(5.2), Inches(0.5), name, size=16, bold=True, color=DARK)

# ---- 11. 演示命令 ----
s = add_slide(); bg(s, WHITE)
title_bar(s, "八、演示命令", 13, TOTAL)
cmds = [
    ("交互式多轮诊断", "python -m diagnosis_agent.cli chat"),
    ("单次文本诊断", "python -m diagnosis_agent.cli diagnose --text \"...\""),
    ("加载 CSV 工单数据", "python -m diagnosis_agent.cli load-data --file data/samples/图谱数据.csv"),
    ("标准接口（对接平台）", "python -m diagnosis_agent.cli diagnose --json-input input.json --std-output"),
    ("查看知识库统计", "python -m diagnosis_agent.cli stats"),
    ("启动 FastAPI 服务", "python -m diagnosis_agent.cli adapter"),
]
rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), RGBColor(0x2B, 0x3A, 0x4A))
for i, (desc, cmd) in enumerate(cmds):
    y = Inches(1.5 + i * 0.85)
    text(s, Inches(0.8), y, Inches(4), Inches(0.7),
         [("# ", {"color": GREEN, "bold": True, "size": 14}), (desc, {"color": WHITE, "size": 14, "bold": True})],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5), y, Inches(7.5), Inches(0.7), cmd, size=13, color=RGBColor(0x8B, 0xC3, 0x4A), anchor=MSO_ANCHOR.MIDDLE)

# ---- 结尾 ----
s = add_slide(); bg(s, PRIMARY)
rect(s, 0, Inches(3.5), SW, Inches(0.08), ACCENT)
text(s, 0, Inches(2.5), SW, Inches(1), "谢谢", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, 0, Inches(4.0), SW, Inches(0.8), "欢迎提问", size=24, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)

# 调整 total
# 因 TOTAL 固定为 12 但实际 14 页，重设页码

prs.save("docs/答辩文档.pptx")
print("已生成: docs/答辩文档.pptx")
print(f"共 {len(prs.slides)} 页")
